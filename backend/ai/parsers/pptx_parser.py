from pptx import Presentation

class PPTXParser:

    @staticmethod
    def extract_text(file_path: str):

        presentation = Presentation(file_path)

        text_runs = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    text_runs.append(shape.text)

        return "\n".join(text_runs)